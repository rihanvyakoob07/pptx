import os
import copy
import requests
from uuid import uuid4
from io import BytesIO
from pptx import Presentation
from app.core.logger import logger
from app.core.utils.persist_helpers import upload_file
from pptx.shapes.picture import Picture
from pptx.util import Inches  # Added for dimension conversion
from pptx.enum.shapes import MSO_SHAPE_TYPE
def delete_file(file_path):
    """Improved file deletion with path validation"""
    try:
        if os.path.exists(file_path):
            os.remove(file_path)
            logger.info(f"Deleted file: {file_path}")
    except Exception as ex:
        logger.error(f"Error deleting {file_path}: {str(ex)}")

def download_slide(url):
    """Improved download with timeout and error handling"""
    file_name = f"{uuid4()}.pptx"
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        with open(file_name, 'wb') as file_obj:
            file_obj.write(response.content)
        logger.info(f"Downloaded: {file_name}")
        return file_name
    except Exception as e:
        logger.error(f"Download failed for {url}: {str(e)}")
        return None

import tempfile
from PIL import Image

def copy_slide_from_external_prs(prs, source_file):
    """Enhanced slide copying with image format validation, temporary file handling, and improved logging"""
    try:
        external_prs = Presentation(source_file)
        for slide_number, source_slide in enumerate(external_prs.slides, start=1):
            logger.info(f"Copying slide {slide_number} from {source_file}")
            
            slide_layout = prs.slide_layouts[6]
            new_slide = prs.slides.add_slide(slide_layout)
            logger.info("Added new slide with blank layout")

            for shape in source_slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        # Validate and convert image format using PIL
                        img = Image.open(BytesIO(shape.image.blob))
                        if img.format not in ['PNG', 'JPEG']:
                            logger.warning(f"Unsupported image format: {img.format}. Converting to PNG.")
                            with BytesIO() as output:
                                img.save(output, format="PNG")
                                image_stream = BytesIO(output.getvalue())
                        else:
                            image_stream = BytesIO(shape.image.blob)
                        
                        # Use temporary file to add picture
                        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{img.format.lower()}") as tmp_img:
                            tmp_img.write(image_stream.getvalue())
                            tmp_img_path = tmp_img.name
                        
                        new_slide.shapes.add_picture(
                            tmp_img_path,
                            shape.left,
                            shape.top,
                            shape.width,
                            shape.height
                        )
                        logger.info("Image copied successfully")
                        
                        # Remove temporary image file
                        os.remove(tmp_img_path)
                    except Exception as img_error:
                        logger.error(f"Failed to copy image: {str(img_error)}")
                else:
                    try:
                        new_element = copy.deepcopy(shape.element)
                        new_slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')
                        logger.info("Shape copied successfully")
                    except Exception as shape_error:
                        logger.error(f"Failed to copy shape: {str(shape_error)}")
    except Exception as e:
        logger.error(f"Slide copy failed for {source_file}: {str(e)}")
        raise



async def generate_combined_slides(slide_urls):
    temp_files = []
    merged_url = None
    
    try:
        logger.info(f"Starting merge of {len(slide_urls)} slides")
        
        # Download slides
        slides_to_merge = []
        for url in slide_urls:
            if downloaded := download_slide(url):
                temp_files.append(downloaded)
                slides_to_merge.append(downloaded)
        
        if not slides_to_merge:
            raise ValueError("No valid slides downloaded")

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(13.33)  # Standard 16:9 aspect ratio
        prs.slide_height = Inches(7.5)

        # Add slides
        for idx, slide_file in enumerate(slides_to_merge, 1):
            logger.info(f"Processing slide {idx}/{len(slides_to_merge)}")
            try:
                copy_slide_from_external_prs(prs, slide_file)
            except Exception as e:
                logger.error(f"Skipping invalid slide {slide_file}: {str(e)}")
                continue

        # Save presentation
        output_file = f"merged-{uuid4()}.pptx"
        prs.save(output_file)
        temp_files.append(output_file)
        logger.info(f"Presentation saved: {output_file}")

        # Upload presentation
        with open(output_file, "rb") as f:
            file_bytes = BytesIO(f.read())
            merged_url = await upload_file(file_bytes, output_file)
            logger.info(f"Upload successful: {merged_url}")

        return merged_url

    except Exception as e:
        logger.error(f"Merge failed: {str(e)}", exc_info=True)
        raise
    finally:
        # Cleanup temp files after successful upload
        for file in temp_files:
            delete_file(file)
        logger.info("Cleanup completed")

